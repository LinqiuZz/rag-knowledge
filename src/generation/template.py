"""模板解析与填充引擎

新方案核心: 用户上传 Word/PPT 模板，系统自动识别占位符，
根据知识库内容或 LLM 生成的内容填充，输出完整的 Office 文件。

模板设计规范:
  - Word 模板: 书签（Bookmark）定位复杂内容，{{key}} 标记纯文本区域
  - PPT 模板: 占位符索引约定（idx=0 标题, idx=1 正文, idx=2 图片）

工作流:
  1. 用户上传模板 → 系统解析并存储
  2. 用户提出需求："根据Q3销售数据填充此模板"
  3. 可选: RAG 检索知识库，获取相关数据块
  4. 构造 Prompt，将模板占位符清单和参考资料发给 LLM
  5. 校验 JSON 格式，调用 python-docx/python-pptx 填充
  6. 生成文件存入 MinIO，返回下载链接
"""

from __future__ import annotations

import json
import logging
import re
from pathlib import Path
from dataclasses import dataclass, field
from typing import Optional
from io import BytesIO

logger = logging.getLogger(__name__)


# ═══════════════════════════════════════════════════════════════
# 数据结构
# ═══════════════════════════════════════════════════════════════

@dataclass
class Placeholder:
    """模板占位符。"""
    id: str                        # 占位符标识
    type: str                      # text | table | image | list
    location: str                  # 位置描述 (bookmark:xxx | placeholder_idx:N)
    hint: str = ""                 # 填充提示
    default: str = ""              # 默认值


@dataclass
class TemplateSchema:
    """模板描述 JSON。"""
    type: str                      # docx | pptx
    name: str = ""
    placeholders: list[Placeholder] = field(default_factory=list)
    styles: list[str] = field(default_factory=list)  # 可用样式列表
    layouts: list[str] = field(default_factory=list)  # PPT 版式列表


# ═══════════════════════════════════════════════════════════════
# 模板解析
# ═══════════════════════════════════════════════════════════════

def parse_template(file_path: str | Path) -> TemplateSchema:
    """解析模板文件，提取占位符清单。

    Args:
        file_path: 模板文件路径 (.docx 或 .pptx)

    Returns:
        TemplateSchema
    """
    path = Path(file_path)
    suffix = path.suffix.lower()

    if suffix == '.docx':
        return _parse_word_template(path)
    elif suffix == '.pptx':
        return _parse_ppt_template(path)
    else:
        raise ValueError(f"不支持的模板格式: {suffix}，仅支持 .docx 和 .pptx")


def _parse_word_template(path: Path) -> TemplateSchema:
    """解析 Word 模板。"""
    try:
        from docx import Document
    except ImportError:
        raise ImportError("请安装 python-docx")

    doc = Document(str(path))
    placeholders = []

    # ── 书签检测 ──
    # python-docx 对书签的支持有限，通过 XML 遍历
    bookmarks_found = set()
    try:
        from docx.opc.constants import RELATIONSHIP_TYPE as RT
        for para in doc.paragraphs:
            for run in para.runs:
                # 检查书签标记
                bm_start = run._element.findall(
                    './/' + '{http://schemas.openxmlformats.org/wordprocessingml/2006/main}bookmarkStart'
                )
                for bm in bm_start:
                    name = bm.get(
                        '{http://schemas.openxmlformats.org/wordprocessingml/2006/main}name'
                    )
                    if name and name not in bookmarks_found:
                        bookmarks_found.add(name)
                        placeholders.append(Placeholder(
                            id=name,
                            type="text",
                            location=f"bookmark:{name}",
                            hint=f"书签位置: {name}",
                        ))
    except Exception as e:
        logger.debug("书签检测: %s", e)

    # ── {{key}} 标记检测 ──
    marker_pattern = re.compile(r'\{\{(\w+)\}\}')
    for para in doc.paragraphs:
        markers = marker_pattern.findall(para.text)
        for marker in markers:
            if marker not in {p.id for p in placeholders}:
                placeholders.append(Placeholder(
                    id=marker,
                    type="text",
                    location=f"marker:{{{{{marker}}}}}",
                    hint=f"文本标记: {{{{ {marker} }}}}",
                ))

    # ── 表格检测（带书签的表可标记为 table 类型）──
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for para in cell.paragraphs:
                    markers = marker_pattern.findall(para.text)
                    for marker in markers:
                        if marker in {p.id for p in placeholders}:
                            # 更新类型为 table
                            for p in placeholders:
                                if p.id == marker:
                                    p.type = "table"
                                    p.hint += " (表格区域)"

    # ── 样式收集 ──
    styles = []
    for style in doc.styles:
        if style.name and style.name not in styles:
            styles.append(style.name)

    return TemplateSchema(
        type="docx",
        name=path.stem,
        placeholders=placeholders,
        styles=styles,
    )


def _parse_ppt_template(path: Path) -> TemplateSchema:
    """解析 PPT 模板。"""
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError("请安装 python-pptx")

    prs = Presentation(str(path))
    placeholders = []
    layouts = []

    # ── 收集版式 ──
    for layout in prs.slide_layouts:
        layouts.append(layout.name)

    # ── 遍历幻灯片检测占位符 ──
    marker_pattern = re.compile(r'\{\{(\w+)\}\}')

    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.is_placeholder:
                ph = shape.placeholder_format
                placeholder_id = f"slide{slide_idx + 1}_idx{ph.idx}"
                placeholder_type = "text"
                hint = f"幻灯片{slide_idx + 1} 占位符 idx={ph.idx}"

                if ph.idx == 0:
                    placeholder_type = "text"
                    hint = f"幻灯片{slide_idx + 1} 标题"
                elif ph.idx == 1:
                    placeholder_type = "text"
                    hint = f"幻灯片{slide_idx + 1} 正文"

                placeholders.append(Placeholder(
                    id=placeholder_id,
                    type=placeholder_type,
                    location=f"slide:{slide_idx + 1},placeholder_idx:{ph.idx}",
                    hint=hint,
                ))

            # {{key}} 标记检测
            if shape.has_text_frame:
                markers = marker_pattern.findall(shape.text_frame.text)
                for marker in markers:
                    placeholders.append(Placeholder(
                        id=marker,
                        type="text",
                        location=f"slide:{slide_idx + 1},marker:{{{{{marker}}}}}",
                        hint=f"幻灯片{slide_idx + 1} 文本标记: {{{{ {marker} }}}}",
                    ))

    return TemplateSchema(
        type="pptx",
        name=path.stem,
        placeholders=placeholders,
        layouts=layouts,
    )


# ═══════════════════════════════════════════════════════════════
# 内容填充工作流
# ═══════════════════════════════════════════════════════════════

class DocumentGenerator:
    """文档生成器 — 模板填充 + LLM 内容生成。

    用法:
        gen = DocumentGenerator(settings, llm)
        # 解析模板
        schema = gen.parse_template("template.docx")
        # 从知识库获取内容
        rag_context = gen.retrieve_context("Q3销售数据")
        # 生成内容
        content_json = gen.generate_content(schema, "根据Q3销售数据填充", rag_context)
        # 填充模板
        result_path = gen.fill_template("template.docx", content_json, "output.docx")
    """

    def __init__(self, settings, llm=None, storage=None):
        from ..config import Settings
        self.settings: Settings = settings
        self.llm = llm
        self.storage = storage

    def parse_template(self, file_path: str | Path) -> TemplateSchema:
        """解析模板。"""
        return parse_template(file_path)

    def retrieve_context(self, query: str, embedder=None, vector_store=None,
                         meta_store=None, user_id: int = None) -> str:
        """从知识库检索相关上下文。

        Returns:
            格式化的上下文文本
        """
        try:
            from ..query.retrieval import RetrievalPipeline, build_rag_context

            pipeline = RetrievalPipeline(self.settings)
            result = pipeline.retrieve(
                query=query,
                embedder=embedder,
                vector_store=vector_store,
                meta_store=meta_store,
                user_id=user_id,
            )
            context, _ = build_rag_context(result.hits)
            return context
        except Exception as e:
            logger.warning("知识库检索失败: %s", e)
            return ""

    def generate_content(
        self,
        schema: TemplateSchema,
        user_request: str,
        rag_context: str = "",
    ) -> dict:
        """调用 LLM 生成填充内容 JSON。

        Args:
            schema: 模板占位符描述
            user_request: 用户需求描述
            rag_context: RAG 检索到的知识库上下文

        Returns:
            填充内容 dict {placeholder_id: content}
        """
        if not self.llm:
            raise RuntimeError("LLM 不可用，无法生成内容")

        # 构建 Prompt
        placeholders_desc = "\n".join(
            f"  - {p.id} (类型: {p.type}, 提示: {p.hint or '无'})"
            for p in schema.placeholders
        )

        system_prompt = """你是一个专业的文档内容生成助手。根据模板占位符和参考资料，生成准确、结构化的填充内容。

规则：
1. 严格基于参考资料提供数据，不要杜撰数字
2. 输出必须是合法的 JSON 格式
3. 每个占位符的填充内容要符合其类型（text/table/list）"""

        user_prompt = f"""请根据以下信息生成模板填充内容。

模板类型: {schema.type}
占位符列表:
{placeholders_desc}

用户需求: {user_request}

参考资料:
{rag_context or "（无参考资料，请根据通用知识填充）"}

请输出 JSON 格式，键为占位符 ID，值为填充内容:
```json
{{
  "placeholder_id1": "填充文本...",
  "placeholder_id2": "填充文本...",
  ...
}}
```

对于 table 类型的占位符，值应为二维数组:
  "sales_table": [["列1", "列2"], ["数据1", "数据2"]]

对于 list 类型的占位符，值应为数组:
  "items": ["项目1", "项目2", "项目3"]
"""

        try:
            response = self.llm.chat(system_prompt, user_prompt, max_tokens=4096)

            # 提取 JSON
            json_match = re.search(r'```(?:json)?\s*([\s\S]*?)```', response)
            if json_match:
                content = json.loads(json_match.group(1))
            else:
                # 尝试直接解析
                content = json.loads(response)

            return content
        except json.JSONDecodeError as e:
            logger.error("LLM 返回的内容不是有效 JSON: %s", e)
            raise ValueError(f"LLM 生成的内容格式有误: {e}")
        except Exception as e:
            logger.error("内容生成失败: %s", e)
            raise

    def fill_template(
        self,
        template_path: str | Path,
        content: dict,
        output_path: str | Path = None,
    ) -> Path:
        """用内容填充模板，生成最终文件。

        Args:
            template_path: 模板文件路径
            content: 填充内容 dict
            output_path: 输出路径（默认在同目录生成）

        Returns:
            输出文件路径
        """
        path = Path(template_path)
        suffix = path.suffix.lower()

        if output_path is None:
            output_path = path.parent / f"{path.stem}_filled{suffix}"
        else:
            output_path = Path(output_path)

        if suffix == '.docx':
            return _fill_word(path, content, output_path)
        elif suffix == '.pptx':
            return _fill_ppt(path, content, output_path)
        else:
            raise ValueError(f"不支持的模板格式: {suffix}")

    def fill_and_upload(
        self,
        template_path: str | Path,
        content: dict,
        object_name: str = None,
    ) -> str:
        """填充模板并上传到 MinIO，返回下载链接。"""
        output_path = self.fill_template(template_path, content)

        if self.storage:
            obj_name = object_name or output_path.name
            self.storage.save_generated(
                output_path.read_bytes(),
                obj_name,
                content_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                if output_path.suffix == '.docx'
                else "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )
            return self.storage.get_generated_url(obj_name)

        return str(output_path)


# ═══════════════════════════════════════════════════════════════
# Word 填充
# ═══════════════════════════════════════════════════════════════

def _fill_word(template_path: Path, content: dict, output_path: Path) -> Path:
    """填充 Word 模板。"""
    try:
        from docx import Document
    except ImportError:
        raise ImportError("请安装 python-docx")

    doc = Document(str(template_path))

    # ── 书签替换 ──
    _fill_word_bookmarks(doc, content)

    # ── {{key}} 文本标记替换 ──
    _fill_word_markers(doc, content)

    # ── 表格填充（如果 content 中有 table 类型数据）──
    _fill_word_tables(doc, content)

    # 保存
    output_path.parent.mkdir(parents=True, exist_ok=True)
    doc.save(str(output_path))
    logger.info("Word 模板填充完成: %s", output_path)
    return output_path


def _fill_word_bookmarks(doc, content: dict):
    """替换 Word 书签。"""
    # python-docx 对书签的处理有限，通过 XML 操作
    try:
        from lxml import etree

        nsmap = {
            'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main',
        }

        for para in doc.paragraphs:
            for run in para.runs:
                bm_starts = run._element.findall('.//w:bookmarkStart', nsmap)
                for bm in bm_starts:
                    name = bm.get('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}name')
                    if name and name in content:
                        value = content[name]
                        if isinstance(value, list):
                            value = "\n".join(str(v) for v in value)
                        # 替换书签后的文本
                        run.text = str(value)
    except ImportError:
        logger.debug("lxml 未安装，跳过书签替换")
    except Exception as e:
        logger.warning("书签替换失败: %s", e)


def _fill_word_markers(doc, content: dict):
    """替换 {{key}} 标记。"""
    marker_pattern = re.compile(r'\{\{(\w+)\}\}')

    for para in doc.paragraphs:
        if marker_pattern.search(para.text):
            text = para.text
            for key, value in content.items():
                if not isinstance(value, (str, int, float)):
                    continue
                pattern = f"{{{{{key}}}}}"
                if pattern in text:
                    # 保留格式的替换
                    for run in para.runs:
                        if pattern in run.text:
                            run.text = run.text.replace(pattern, str(value))
                    text = para.text  # 更新文本

    # 也处理表格中的标记
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for para in cell.paragraphs:
                    if marker_pattern.search(para.text):
                        for run in para.runs:
                            for key, value in content.items():
                                if isinstance(value, (str, int, float)):
                                    pattern = "{{" + key + "}}"
                                    if pattern in run.text:
                                        run.text = run.text.replace(pattern, str(value))


def _fill_word_tables(doc, content: dict):
    """填充 Word 表格。"""
    for key, value in content.items():
        if not isinstance(value, list):
            continue
        if not value or not isinstance(value[0], list):
            continue

        # 查找包含 {{key}} 标记的表格
        target_table = None
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if f"{{{{{key}}}}}" in cell.text:
                        target_table = table
                        break
                if target_table:
                    break
            if target_table:
                break

        if target_table:
            # 清空并重建表格
            _rebuild_table(target_table, value)


def _rebuild_table(table, data: list[list]):
    """重建表格内容。"""
    # 保留表头行，清除数据行
    while len(table.rows) > 1:
        tr = table.rows[1]._tr
        table._tbl.remove(tr)

    # 填充数据
    for row_data in data[1:]:  # 跳过表头
        row = table.add_row()
        for i, cell_text in enumerate(row_data):
            if i < len(row.cells):
                row.cells[i].text = str(cell_text)


# ═══════════════════════════════════════════════════════════════
# PPT 填充
# ═══════════════════════════════════════════════════════════════

def _fill_ppt(template_path: Path, content: dict, output_path: Path) -> Path:
    """填充 PPT 模板。"""
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError("请安装 python-pptx")

    prs = Presentation(str(template_path))

    marker_pattern = re.compile(r'\{\{(\w+)\}\}')

    for slide in prs.slides:
        for shape in slide.shapes:
            # ── 占位符按索引填充 ──
            if shape.is_placeholder:
                idx = shape.placeholder_format.idx
                # 查找匹配的 content key
                for key, value in content.items():
                    if f"idx{idx}" in key or f"idx_{idx}" in key:
                        if shape.has_text_frame:
                            shape.text_frame.text = str(value)
                        break

            # ── {{key}} 标记替换 ──
            if shape.has_text_frame:
                markers = marker_pattern.findall(shape.text_frame.text)
                for marker in markers:
                    if marker in content:
                        value = content[marker]
                        if isinstance(value, str):
                            shape.text_frame.text = shape.text_frame.text.replace(
                                f"{{{{{marker}}}}}", value
                            )

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    logger.info("PPT 模板填充完成: %s", output_path)
    return output_path
