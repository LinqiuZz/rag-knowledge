"""多格式文档深度解析器 — 保留完整排版信息

新方案核心: 将二进制文件转为富含语义标签的结构化 JSON，
同时输出纯文本供向量化。

支持的格式:
  - Word (.docx): 段落样式、表格、页眉页脚、脚注尾注、图片替代文本
  - PPT (.pptx): 幻灯片标题、布局、文本框位置、演讲者备注、图表
  - Excel (.xlsx): 工作表、行列数据
  - PDF: 数字原生 + OCR 扫描件
  - 图片: OCR + 多模态描述
  - Markdown / TXT / HTML / Code: 纯文本解析

输出: DocumentStructure 包含 page/slide 级元数据和 element 列表
"""

from __future__ import annotations

import re
import logging
from pathlib import Path
from dataclasses import dataclass, field
from typing import Optional

logger = logging.getLogger(__name__)


# ═══════════════════════════════════════════════════════════════
# 数据结构
# ═══════════════════════════════════════════════════════════════

@dataclass
class PageElement:
    """页面/幻灯片中的单个元素。"""
    type: str                      # Heading1-9, Paragraph, Table, Picture, TextBox, Header, Footer, ListItem, Caption, Chart
    text: str
    style: str = ""                # 样式名 (Normal, Heading 1, etc.)
    level: int = 0                 # 标题层级
    position: str = ""             # left/right/top/bottom/full — 用于 PPT 布局
    bbox: tuple = None             # 边界框 (x0, y0, x1, y1) — PDF/PPT
    metadata: dict = field(default_factory=dict)  # 额外信息（表格行列数、图片描述等）


@dataclass
class PageData:
    """单页/单张幻灯片数据。"""
    page_number: int
    header: str = ""
    footer: str = ""
    elements: list[PageElement] = field(default_factory=list)
    raw_text: str = ""             # 纯文本（用于向量化）


@dataclass
class SlideData:
    """单张幻灯片数据（PPT 专用）。"""
    slide_number: int
    layout: str = ""               # 布局名称（如"两栏内容"）
    title: str = ""
    notes: str = ""                # 演讲者备注
    elements: list[PageElement] = field(default_factory=list)
    raw_text: str = ""


@dataclass
class DocumentStructure:
    """完整的文档结构化表示。"""
    doc_id: str = ""
    title: str = ""
    format: str = ""               # docx, pptx, xlsx, pdf, image, md, txt, html, code
    # 文档级元数据
    author: str = ""
    created_at: str = ""
    modified_at: str = ""
    page_count: int = 0
    char_count: int = 0
    # 内容
    pages: list[PageData] = field(default_factory=list)
    slides: list[SlideData] = field(default_factory=list)
    # 纯文本（所有元素文本的拼接）
    full_text: str = ""


# ═══════════════════════════════════════════════════════════════
# 解析器注册表
# ═══════════════════════════════════════════════════════════════

SUPPORTED_EXTENSIONS = {
    # Office
    '.docx': 'word',
    '.pptx': 'ppt',
    '.xlsx': 'excel',
    '.xls': 'excel',
    # 文档
    '.pdf': 'pdf',
    '.md': 'text',
    '.markdown': 'text',
    '.txt': 'text',
    '.text': 'text',
    '.html': 'text',
    '.htm': 'text',
    # 图片
    '.png': 'image',
    '.jpg': 'image',
    '.jpeg': 'image',
    '.gif': 'image',
    '.bmp': 'image',
    '.tiff': 'image',
    # 配置/数据
    '.json': 'text',
    '.yaml': 'text',
    '.yml': 'text',
    '.toml': 'text',
    '.ini': 'text',
    '.cfg': 'text',
    '.conf': 'text',
    '.xml': 'text',
    '.csv': 'text',
    # 代码
    '.py': 'code', '.js': 'code', '.ts': 'code', '.java': 'code',
    '.cpp': 'code', '.c': 'code', '.h': 'code', '.go': 'code',
    '.rs': 'code', '.rb': 'code', '.php': 'code', '.swift': 'code',
    '.kt': 'code', '.scala': 'code', '.sql': 'code',
    '.sh': 'code', '.bash': 'code', '.bat': 'code',
}


def parse_document(file_path: str | Path) -> DocumentStructure:
    """自动识别并深度解析文档。

    Args:
        file_path: 文件路径

    Returns:
        携带完整排版信息的 DocumentStructure
    """
    path = Path(file_path)
    if not path.exists():
        raise FileNotFoundError(f"文件不存在: {path}")

    suffix = path.suffix.lower()
    doc_type = SUPPORTED_EXTENSIONS.get(suffix)

    if doc_type is None:
        if suffix == ".doc":
            raise ValueError("不支持旧版 .doc 格式，请转换为 .docx 格式后重试")
        raise ValueError(f"不支持的文件格式: {suffix}")

    parser = _PARSER_MAP.get(doc_type)
    if parser is None:
        raise ValueError(f"解析器未实现: {doc_type}")

    return parser(path)


# ═══════════════════════════════════════════════════════════════
# Word (.docx) 深度解析
# ═══════════════════════════════════════════════════════════════

def parse_word_deep(path: Path) -> DocumentStructure:
    """深度解析 Word 文档，保留段落样式、表格、页眉页脚。"""
    try:
        from docx import Document
        from docx.enum.style import WD_STYLE_TYPE
    except ImportError:
        raise ImportError("请安装 python-docx: pip install python-docx")

    doc = Document(str(path))

    # 提取所有样式
    style_map = {}
    for style in doc.styles:
        if style.type == WD_STYLE_TYPE.PARAGRAPH:
            style_map[style.style_id] = style.name

    pages = []
    all_text_parts = []

    # 遍历段落
    current_page_elements = []
    current_page_num = 1

    for para in doc.paragraphs:
        if not para.text.strip() and not _has_any_runs(para):
            # 空段落可能是分页标记
            if _is_page_break(para):
                if current_page_elements:
                    pages.append(PageData(
                        page_number=current_page_num,
                        elements=current_page_elements,
                        raw_text="\n".join(e.text for e in current_page_elements),
                    ))
                    current_page_elements = []
                    current_page_num += 1
            continue

        text = para.text.strip()
        style_name = para.style.name if para.style else "Normal"

        # 判断元素类型
        elem_type = _word_element_type(style_name, text)

        element = PageElement(
            type=elem_type,
            text=text,
            style=style_name,
            level=_heading_level(style_name),
            metadata={"alignment": str(para.alignment) if para.alignment else None},
        )
        current_page_elements.append(element)
        all_text_parts.append(text)

    # 提取表格
    for table in doc.tables:
        rows_data = []
        for row in table.rows:
            row_cells = [cell.text.strip() for cell in row.cells]
            rows_data.append(row_cells)

        # 将表格转为 Markdown 表格文本
        table_text = _table_to_text(rows_data)
        element = PageElement(
            type="Table",
            text=table_text,
            style="Table Grid",
            metadata={"rows": len(rows_data), "cols": len(rows_data[0]) if rows_data else 0},
        )
        current_page_elements.append(element)
        all_text_parts.append(table_text)

    # 最后一页
    if current_page_elements:
        pages.append(PageData(
            page_number=current_page_num,
            elements=current_page_elements,
            raw_text="\n".join(e.text for e in current_page_elements),
        ))

    # 文档属性
    props = doc.core_properties
    full_text = "\n\n".join(all_text_parts)

    return DocumentStructure(
        doc_id=path.stem,
        title=path.stem,
        format="docx",
        author=str(props.author) if props.author else "",
        created_at=str(props.created) if props.created else "",
        modified_at=str(props.modified) if props.modified else "",
        page_count=len(pages),
        char_count=len(full_text),
        pages=pages,
        full_text=full_text,
    )


def _word_element_type(style_name: str, text: str) -> str:
    """根据样式名推断 Word 元素类型。"""
    style_lower = style_name.lower()
    if any(kw in style_lower for kw in ('heading', '标题')):
        level = _heading_level(style_name)
        return f"Heading{level}" if level else "Heading"
    if 'list' in style_lower or 'bullet' in style_lower:
        return "ListItem"
    if 'caption' in style_lower or '题注' in style_lower:
        return "Caption"
    if 'header' in style_lower or '页眉' in style_lower:
        return "Header"
    if 'footer' in style_lower or '页脚' in style_lower:
        return "Footer"
    return "Paragraph"


def _heading_level(style_name: str) -> int:
    """从样式名提取标题层级。"""
    match = re.search(r'(\d+)', style_name)
    if match:
        return int(match.group(1))
    # 中文标题样式
    cn_map = {
        '标题': 1, '标题 1': 1, '标题1': 1,
        '标题 2': 2, '标题2': 2, '副标题': 2,
        '标题 3': 3, '标题3': 3,
        '标题 4': 4, '标题4': 4,
    }
    return cn_map.get(style_name.strip(), 0)


def _has_any_runs(para) -> bool:
    """段落是否有任何 run 内容（文字或图片）。"""
    try:
        return any(r.text or r._element.findall('.//' + '{http://schemas.openxmlformats.org/wordprocessingml/2006/main}drawing') for r in para.runs)
    except Exception:
        return bool(para.runs)


def _is_page_break(para) -> bool:
    """检测段落是否包含分页符。"""
    try:
        for run in para.runs:
            if run._element.findall('.//' + '{http://schemas.openxmlformats.org/wordprocessingml/2006/main}br'):
                return True
    except Exception:
        pass
    return False


def _table_to_text(rows: list[list[str]]) -> str:
    """将表格数据转为可读文本。"""
    if not rows:
        return ""
    lines = []
    # 表头
    lines.append("| " + " | ".join(rows[0]) + " |")
    lines.append("| " + " | ".join(["---" for _ in rows[0]]) + " |")
    # 数据行
    for row in rows[1:]:
        padded = row + [""] * (len(rows[0]) - len(row))
        lines.append("| " + " | ".join(padded[:len(rows[0])]) + " |")
    return "\n".join(lines)


# ═══════════════════════════════════════════════════════════════
# PPT (.pptx) 深度解析
# ═══════════════════════════════════════════════════════════════

def parse_ppt_deep(path: Path) -> DocumentStructure:
    """深度解析 PPT，保留布局、位置、备注等信息。"""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        raise ImportError("请安装 python-pptx: pip install python-pptx")

    prs = Presentation(str(path))
    slides = []
    all_text_parts = []

    for i, slide in enumerate(prs.slides, 1):
        # 布局名称
        layout_name = slide.slide_layout.name if slide.slide_layout else ""

        # 标题
        title_text = ""
        if slide.shapes.title:
            title_text = slide.shapes.title.text.strip()

        # 备注
        notes_text = ""
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()

        elements = []

        for shape in slide.shapes:
            elem = _parse_ppt_shape(shape)
            if elem:
                elements.append(elem)

        # 合成纯文本
        raw_text = title_text
        if elements:
            raw_text += "\n" + "\n".join(e.text for e in elements if e.type != "Picture")
        if notes_text:
            raw_text += f"\n[备注] {notes_text}"

        slides.append(SlideData(
            slide_number=i,
            layout=layout_name,
            title=title_text,
            notes=notes_text,
            elements=elements,
            raw_text=raw_text.strip(),
        ))
        all_text_parts.append(raw_text.strip())

    full_text = "\n\n---\n\n".join(all_text_parts)

    return DocumentStructure(
        doc_id=path.stem,
        title=path.stem,
        format="pptx",
        page_count=len(slides),
        char_count=len(full_text),
        slides=slides,
        full_text=full_text,
    )


def _parse_ppt_shape(shape) -> PageElement | None:
    """解析 PPT 中的单个形状。"""
    shape_type = str(shape.shape_type) if shape.shape_type else "UNKNOWN"

    # 组合（GroupShapes）
    if shape_type == "GROUP (6)":
        texts = []
        try:
            for child in shape.shapes:
                if child.has_text_frame:
                    texts.append(child.text_frame.text.strip())
        except Exception:
            pass
        if texts:
            return PageElement(type="TextBox", text="\n".join(texts), position="full")
        return None

    # 文本框 / 占位符
    if shape.has_text_frame:
        text = shape.text_frame.text.strip()
        if not text:
            return None

        # 判断是否为占位符
        is_placeholder = shape.is_placeholder
        placeholder_idx = shape.placeholder_format.idx if is_placeholder else None

        # 位置信息
        left = shape.left
        top = shape.top
        width = shape.width
        height = shape.height

        position = _ppt_position(left, top, width, height)

        elem_type = "TextBox"
        if is_placeholder and placeholder_idx == 0:
            elem_type = "Title"
        elif is_placeholder and placeholder_idx == 1:
            elem_type = "Body"

        return PageElement(
            type=elem_type,
            text=text,
            position=position,
            metadata={
                "placeholder_idx": placeholder_idx,
                "is_placeholder": is_placeholder,
            },
        )

    # 表格
    if shape.has_table:
        table = shape.table
        rows = []
        for row in table.rows:
            rows.append([cell.text.strip() for cell in row.cells])
        table_text = _table_to_text(rows)
        return PageElement(
            type="Table",
            text=table_text,
            metadata={"rows": len(rows), "cols": len(rows[0]) if rows else 0},
        )

    # 图片
    if shape_type in ("PICTURE (13)", "13"):
        return PageElement(
            type="Picture",
            text="",
            metadata={"image_type": "embedded", "description": "[待多模态模型生成]"},
        )

    # 图表
    if shape.has_chart:
        return PageElement(
            type="Chart",
            text="[图表数据]",
            metadata={"chart_type": str(shape.chart.chart_type)},
        )

    return None


def _ppt_position(left, top, width, height) -> str:
    """推断元素在幻灯片上的位置（left/right/top/bottom/full）。"""
    try:
        from pptx.util import Inches
        # 大约位置判断（基于 10 英寸宽标准幻灯片）
        w_inches = width / 914400  # EMU to inches
        l_inches = left / 914400
        if w_inches > 8 and l_inches < 1:
            return "full"
        if l_inches < 4:
            return "left"
        else:
            return "right"
    except Exception:
        return ""


# ═══════════════════════════════════════════════════════════════
# PDF 深度解析
# ═══════════════════════════════════════════════════════════════

def parse_pdf_deep(path: Path) -> DocumentStructure:
    """深度解析 PDF，提取每页文本和坐标信息。"""
    try:
        import fitz  # PyMuPDF
    except ImportError:
        raise ImportError("请安装 PyMuPDF: pip install PyMuPDF")

    doc = fitz.open(str(path))
    pages = []
    all_text_parts = []

    for i in range(len(doc)):
        page = doc[i]
        # 提取带位置信息的文本块
        blocks = page.get_text("blocks")
        elements = []
        page_text_parts = []

        for block in blocks:
            # block = (x0, y0, x1, y1, text, block_no, block_type)
            if len(block) >= 5:
                text = block[4].strip() if isinstance(block[4], str) else ""
                bbox = tuple(block[:4])
                if text:
                    elements.append(PageElement(
                        type="Paragraph",
                        text=text,
                        bbox=bbox,
                    ))
                    page_text_parts.append(text)

        raw_text = "\n".join(page_text_parts)
        pages.append(PageData(
            page_number=i + 1,
            elements=elements,
            raw_text=raw_text,
        ))
        all_text_parts.append(raw_text)

    doc.close()
    full_text = "\n\n".join(all_text_parts)

    return DocumentStructure(
        doc_id=path.stem,
        title=path.stem,
        format="pdf",
        page_count=len(pages),
        char_count=len(full_text),
        pages=pages,
        full_text=full_text,
    )


# ═══════════════════════════════════════════════════════════════
# Excel 解析
# ═══════════════════════════════════════════════════════════════

def parse_excel(path: Path) -> DocumentStructure:
    """解析 Excel 文档。"""
    try:
        import openpyxl
    except ImportError:
        raise ImportError("请安装 openpyxl: pip install openpyxl")

    wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    pages = []
    all_text_parts = []

    for idx, sheet_name in enumerate(wb.sheetnames, 1):
        ws = wb[sheet_name]
        rows_data = []
        for row in ws.iter_rows(values_only=True):
            row_text = [str(cell) if cell is not None else "" for cell in row]
            if any(row_text):
                rows_data.append(row_text)

        table_text = _table_to_text(rows_data) if rows_data else ""
        sheet_text = f"【{sheet_name}】\n{table_text}"

        pages.append(PageData(
            page_number=idx,
            elements=[PageElement(
                type="Table",
                text=table_text,
                style=sheet_name,
                metadata={"sheet": sheet_name, "rows": len(rows_data)},
            )],
            raw_text=sheet_text,
        ))
        all_text_parts.append(sheet_text)

    wb.close()
    full_text = "\n\n".join(all_text_parts)

    return DocumentStructure(
        doc_id=path.stem,
        title=path.stem,
        format="xlsx",
        page_count=len(pages),
        char_count=len(full_text),
        pages=pages,
        full_text=full_text,
    )


# ═══════════════════════════════════════════════════════════════
# 图片解析（OCR + 多模态描述）
# ═══════════════════════════════════════════════════════════════

def parse_image(path: Path) -> DocumentStructure:
    """解析图片 — OCR 提取文字 + 标记待多模态描述。

    注：OCR 和多模态描述在生产环境中异步执行，
    此函数返回基础结构，后续异步补充。
    """
    text = f"[图片: {path.name}]"
    # 尝试 OCR（如果 PaddleOCR 可用）
    try:
        from paddleocr import PaddleOCR
        ocr = PaddleOCR(lang='ch', show_log=False)
        result = ocr.ocr(str(path))
        if result and result[0]:
            lines = [line[1][0] for line in result[0]]
            text = "\n".join(lines)
    except ImportError:
        logger.debug("PaddleOCR 未安装，跳过 OCR")
    except Exception as e:
        logger.warning("OCR 失败: %s", e)

    return DocumentStructure(
        doc_id=path.stem,
        title=path.stem,
        format="image",
        page_count=1,
        char_count=len(text),
        pages=[PageData(
            page_number=1,
            elements=[PageElement(type="Picture", text=text, metadata={
                "ocr_text": text,
                "multimodal_description": "[待异步生成]",
            })],
            raw_text=text,
        )],
        full_text=text,
    )


# ═══════════════════════════════════════════════════════════════
# 文本/代码解析
# ═══════════════════════════════════════════════════════════════

def parse_text(path: Path) -> DocumentStructure:
    """解析文本文件（Markdown、TXT、HTML、JSON等）。"""
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        text = f.read()

    if not text.strip():
        raise ValueError("文件内容为空")

    # Markdown 特殊处理：按标题分段
    if path.suffix.lower() in ('.md', '.markdown'):
        return _parse_markdown_structured(path, text)

    # HTML 处理
    if path.suffix.lower() in ('.html', '.htm'):
        try:
            from bs4 import BeautifulSoup
            soup = BeautifulSoup(text, "html.parser")
            for tag in soup(["script", "style"]):
                tag.decompose()
            text = soup.get_text(separator="\n", strip=True)
        except ImportError:
            pass

    title = text.strip().split('\n')[0].lstrip('#').strip()[:200] or path.stem

    return DocumentStructure(
        doc_id=path.stem,
        title=title,
        format=path.suffix.lstrip('.'),
        page_count=1,
        char_count=len(text),
        pages=[PageData(
            page_number=1,
            elements=[PageElement(type="Paragraph", text=text)],
            raw_text=text,
        )],
        full_text=text,
    )


def _parse_markdown_structured(path: Path, text: str) -> DocumentStructure:
    """结构化解析 Markdown，保留标题层级。"""
    lines = text.split('\n')
    title = path.stem
    elements = []
    current_text = []
    current_heading = None

    for line in lines:
        # 标题检测
        heading_match = re.match(r'^(#{1,6})\s+(.+)', line)
        if heading_match:
            if current_text:
                elements.append(PageElement(
                    type=f"Heading{current_heading}" if current_heading else "Paragraph",
                    text='\n'.join(current_text),
                    level=current_heading or 0,
                ))
                current_text = []
            level = len(heading_match.group(1))
            heading_text = heading_match.group(2)
            if not title or title == path.stem:
                title = heading_text
            elements.append(PageElement(
                type=f"Heading{level}",
                text=heading_text,
                level=level,
            ))
            current_heading = level
        else:
            if line.strip():
                current_text.append(line)

    # 剩余文本
    if current_text:
        elements.append(PageElement(
            type="Paragraph",
            text='\n'.join(current_text),
        ))

    return DocumentStructure(
        doc_id=path.stem,
        title=title,
        format="md",
        page_count=1,
        char_count=len(text),
        pages=[PageData(
            page_number=1,
            elements=elements,
            raw_text=text,
        )],
        full_text=text,
    )


def parse_code(path: Path) -> DocumentStructure:
    """解析代码文件。"""
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        text = f.read()

    if not text.strip():
        raise ValueError("文件内容为空")

    ext = path.suffix.lstrip('.')
    title = path.stem

    return DocumentStructure(
        doc_id=path.stem,
        title=title,
        format=ext,
        page_count=1,
        char_count=len(text),
        pages=[PageData(
            page_number=1,
            elements=[PageElement(
                type="Code",
                text=text,
                metadata={"language": ext},
                style=f"language-{ext}",
            )],
            raw_text=text,
        )],
        full_text=text,
    )


# ═══════════════════════════════════════════════════════════════
# 解析器注册表（所有解析函数定义完毕后注册）
# ═══════════════════════════════════════════════════════════════

_PARSER_MAP = {
    "word": parse_word_deep,
    "ppt": parse_ppt_deep,
    "excel": parse_excel,
    "pdf": parse_pdf_deep,
    "image": parse_image,
    "text": parse_text,
    "code": parse_code,
}


# ═══════════════════════════════════════════════════════════════
# 工具函数：从 DocumentStructure 提取纯文本列表（用于分块）
# ═══════════════════════════════════════════════════════════════

def extract_text_segments(doc: DocumentStructure) -> list[dict]:
    """从 DocumentStructure 提取带元数据的文本段。

    每个段包含:
      - text: 纯文本
      - page_number / slide_number
      - element_type
      - style / level / position
      - metadata (header, footer, layout, notes 等)
    """
    segments = []

    # Word/PDF 页面
    for page in doc.pages:
        for elem in page.elements:
            segments.append({
                "text": elem.text,
                "page_number": page.page_number,
                "element_type": elem.type,
                "style": elem.style,
                "level": elem.level,
                "header": page.header,
                "footer": page.footer,
                "metadata": elem.metadata,
            })

    # PPT 幻灯片
    for slide in doc.slides:
        # 标题
        if slide.title:
            segments.append({
                "text": slide.title,
                "slide_number": slide.slide_number,
                "element_type": "Title",
                "layout": slide.layout,
                "notes": slide.notes,
            })
        # 元素
        for elem in slide.elements:
            segments.append({
                "text": elem.text,
                "slide_number": slide.slide_number,
                "element_type": elem.type,
                "position": elem.position,
                "layout": slide.layout,
                "notes": slide.notes,
                "metadata": elem.metadata,
            })

    # 如果没有任何页面/幻灯片（纯文本），返回整个文档
    if not segments and doc.full_text:
        segments.append({
            "text": doc.full_text,
            "page_number": 1,
            "element_type": "Paragraph",
        })

    return segments
